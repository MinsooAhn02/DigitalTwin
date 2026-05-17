"""
Traffic Digital Twin - Poster Generator (FULL CONTENT)
Run: python create_poster.py
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from lxml import etree

NAVY   = RGBColor(0x00, 0x20, 0x60)
WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
DARK   = RGBColor(0x1F, 0x2D, 0x3D)
LIGHT  = RGBColor(0xE8, 0xF0, 0xFE)
ACCENT = RGBColor(0x00, 0x70, 0xC0)

W  = 33.11
H  = 46.81
M  = 1.26        # margin
CW = 14.46       # column width
GAP = W - 2*M - 2*CW   # ~1.67"
CR  = M + CW + GAP      # right col x ~17.39"
SH  = 1.16       # section header height
CT  = 8.36       # content top

def inch(x): return Inches(x)

def sec_hdr(slide, text, left, top, width):
    box = slide.shapes.add_shape(1, inch(left), inch(top), inch(width), inch(SH))
    box.fill.solid(); box.fill.fore_color.rgb = NAVY
    box.line.fill.background()
    tf = box.text_frame; tf.word_wrap = False
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "  " + text
    r.font.bold = True; r.font.size = Pt(32); r.font.color.rgb = WHITE
    return top + SH

def txtbox(slide, text, left, top, width, height, size=20, bold=False, color=None):
    color = color or DARK
    tb = slide.shapes.add_textbox(inch(left), inch(top), inch(width), inch(height))
    tf = tb.text_frame; tf.word_wrap = True
    for i, line in enumerate(text.split('\n')):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        r = p.add_run(); r.text = line
        r.font.size = Pt(size); r.font.bold = bold; r.font.color.rgb = color
        pPr = p._p.get_or_add_pPr()
        ls = etree.SubElement(pPr, qn('a:lnSpc'))
        sp = etree.SubElement(ls, qn('a:spcPct')); sp.set('val', '115000')

def img_placeholder(slide, label, left, top, width, height):
    box = slide.shapes.add_shape(1, inch(left), inch(top), inch(width), inch(height))
    box.fill.solid(); box.fill.fore_color.rgb = LIGHT
    box.line.color.rgb = ACCENT; box.line.width = Pt(1.5)
    tf = box.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
    r = p.add_run(); r.text = f"[ INSERT SCREENSHOT ]\n{label}"
    r.font.size = Pt(22); r.font.color.rgb = ACCENT; r.font.bold = True

def build():
    prs = Presentation("Poster Format_blank.poster.pptx")
    slide = prs.slides[0]

    # Title placeholder
    for s in slide.shapes:
        if s.name == "Title 1" and s.shape_type == 14:
            tf = s.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "Traffic Digital Twin"
            r.font.bold = True; r.font.size = Pt(72); r.font.color.rgb = NAVY

    # Authors
    for s in slide.shapes:
        if s.name == "Title 1" and s.shape_type == 17:
            tf = s.text_frame; tf.clear()
            p = tf.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
            r = p.add_run(); r.text = "Minsoo Ahn  |  CSE 327 – Spring 2026  |  Korea University"
            r.font.size = Pt(30); r.font.color.rgb = DARK

    # Subtitle
    txtbox(slide,
        "Real-time Vehicle Detection & Traffic Analytics from Live CCTV Streams",
        2.0, 6.9, 29.0, 1.0, size=28, bold=False)

    # ═══════════════════════════════════════════
    # LEFT COLUMN
    # ═══════════════════════════════════════════
    t = CT

    # 1. Project Overview
    t = sec_hdr(slide, "PROJECT OVERVIEW", M, t, CW)
    txtbox(slide, """A traffic digital twin system that analyzes real road CCTV footage with AI and visualizes vehicle positions, speeds, and traffic flow in real time on an interactive map.

Key Features:
  - Receives live CCTV HLS streams from ITS (National Traffic Information Center) OpenAPI
  - Vehicle detection and multi-object tracking via YOLOv8 + BoxMOT
  - Pixel coordinate → GPS coordinate conversion (4-point homography calibration)
  - Real-time display of vehicle markers, speed, LOS, and alerts on deck.gl map
  - Region of Interest (ROI) editor and camera direction calibration UI""",
        M, t+0.1, CW, 3.5, size=20)
    t += 3.8

    # 2. System Architecture & Data Flow
    t = sec_hdr(slide, "SYSTEM ARCHITECTURE & DATA FLOW", M, t, CW)
    txtbox(slide, """[ Browser → Backend Flow ]

  Click CCTV icon on map
    → POST /switch-camera
    → <video> HLS playback (hls.js)
    → [YOLO tab] <canvas> frame capture (640px JPEG, ~33ms interval)
    → WS /ws/detect ──JPEG──► backend inference ◄── annotated JPEG returned
    → WS /ws ◄── JSON broadcast
       (vehicle positions/speed/stats → map markers + sidebar update)

[ FastAPI Backend — /ws/detect processing order ]
  1. JPEG decode → numpy ndarray
  2. ROI masking (remove detections outside region of interest)
  3. YOLOv8 inference → bounding box list
  4. BoxMOT tracker → assign/maintain track_id
  5. PerspectiveTransformer → pixel→GPS/meter conversion
  6. TrafficAnalytics → speed, LOS, alert calculation
  7. _broadcast() → send JSON to /ws subscribers
  8. Annotated JPEG → return to browser

[ Background Loops ]
  live_loop:        Read HLS stream directly via OpenCV → run same pipeline
                    (continuous server-side processing even without browser)
  hls_refresh_loop: Re-call ITS API every 30 min → refresh expired token → swap URL""",
        M, t+0.1, CW, 8.0, size=19)
    t += 8.3

    # 3. Tech Stack
    t = sec_hdr(slide, "TECH STACK", M, t, CW)
    txtbox(slide, """[ Backend — Python 3.11 ]
  Web server:  FastAPI + Uvicorn (async REST + WebSocket)
  AI detect:   ultralytics YOLOv8x / YOLOv8s
               TensorRT FP16 (.engine) — GPU-accelerated inference
               torch 2.6.0+cu124 (CUDA 12.4)
  Tracking:    BoxMOT — BotSort / ByteTrack / OcSort / DeepOcSort
               (selected via TRACKER_TIER env var, default: auto)
  Video:       OpenCV-Python (cv2), FFmpeg (HLS decoding)
  Coords:      numpy (homography matrix ops)
  HTTP:        httpx (async ITS API calls)
  Config:      python-dotenv, pydantic

[ Frontend — JavaScript (ES2022 + JSX) ]
  Build:  Vite + React 18 (functional components + Hooks)
  Map:    deck.gl (ScatterplotLayer, TextLayer, PathLayer,
                   PolygonLayer, IconLayer)
          react-map-gl + MapLibre GL (vector tiles)
          Carto Dark Matter style theme
  HLS:    hls.js (m3u8 playback in browser)
  Charts: Recharts (pie chart, radial bar)
  Comm:   native browser WebSocket, fetch API
  State:  useState / useReducer / useRef / useCallback / useMemo

[ External Service ]
  ITS OpenAPI: openapi.its.go.kr:9443/cctvInfo
               Query CCTV list (location + HLS URL) within current map viewport bbox
               HLS URL contains expiring token → auto-refresh logic applied""",
        M, t+0.1, CW, 9.5, size=19)
    t += 9.8

    # 4. Backend Modules
    t = sec_hdr(slide, "BACKEND MODULES", M, t, CW)
    txtbox(slide, """[ config.py ]
  YOLO_MODEL, YOLO_CONF(0.25)/YOLO_IOU(0.45), VEHICLE_CLASSES {2:car,3:motorcycle,5:bus,7:truck}
  TRACKER_TIER(auto/cpu/low/medium/high), SPEED_LIMIT_KPH(120), SPEED_SMOOTHING_ALPHA(0.30)
  MAX_REASONABLE_KPH(180), GC_GRACE_FRAMES(5), YOLO_DETECT_INTERVAL(3)
  BOTTLENECK_DWELL_FRAMES(150=5sec), PARKED_FRAMES_THRESHOLD(300=10sec), PARKED_POSITION_RADIUS_PX(30)

[ detector.py — VehicleDetector ]
  - Model loading: auto-select .pt (PyTorch) or .engine (TensorRT FP16)
  - detect(frame): plain YOLO inference (no tracker_id)
  - track(frame): BoxMOT-integrated inference
      1. Run YOLO predict every YOLO_DETECT_INTERVAL(=3) frames;
         reuse previous detections for remaining frames (2/3 GPU load reduction)
      2. half=True (FP16) inference on CUDA (~30-50% speed gain)
      3. ROI masking (apply polygon loaded from roi_manager)
      4. tracker.update(bboxes, frame) → assign/maintain track_id (incl. Kalman prediction)
      5. Return supervision.Detections object
  - set_roi(polygon), reset_tracker(), threading.Lock (serialize live_loop + ws/detect)
  VideoStream: switch_to(url), read_frame(), reconnect()

[ tracker.py — VehicleTracker ]
  - In/Out counting via LineZone crossing detection
  - update(detections, frame_wh): returns supervision LineZone result

[ analytics.py — TrafficAnalytics ]
  Input:  VehicleState list (track_id, class_name, center_px, lat, lon, ...)
  Output: FrameAnalytics (vehicles, vehicle_count, avg_speed_kph, los_grade, ...)
  Processing order:
    1. _gc(active): clean up undetected tracks (delete after 5-frame grace period)
    2. _is_near_parked(): within 30px of registered parked position → is_parked immediately
    3. _speed(): Haversine distance-based speed calculation
       dist<0.20m→speed=0 / raw_kph>180→skip outlier / EMA(alpha=0.15) smoothing
    4. _dwell_update(): accumulate stationary frames → bottleneck/parked detection
    5. Exclude parked vehicles, aggregate stats (avg_speed, LOS, class_counts)

[ transform.py — PerspectiveTransformer ]
  - 4 pairs of (pixel, GPS) correspondences → OpenCV findHomography() → Homography matrix
  - Implicitly encodes camera height, angle, and tilt
  - pixel_to_gps(cx,cy)→(lat,lon), pixel_to_meter(cx,cy)→(x_m,y_m)
  - update_from_calibration(): simultaneously update _H_gps + _H_meter (consistency guaranteed)
  - Before calibration: uses default grid from config.py (inaccurate) / After: precise transform

[ roi_manager.py ]
  - Stored in normalized coordinates ([0,1]) → resolution-independent
  - Saved in roi_config.json keyed by camera URL
  - get_roi / save_roi / delete_roi

[ model_setup.py ]
  - Tkinter GUI: select YOLO model, performance profile (quality/balanced/performance),
    CUDA toggle, tracker tier
  - Saved to .yolo_model + .runtime_profile.json
  - FPS values update live when toggling CUDA checkbox

[ main.py — FastAPI Endpoints ]
  GET  /cctvs           ITS API → CCTV list within viewport (TTL 5min cache)
  POST /switch-camera   switch camera + reset tracker/analytics + send camera_ready signal
  GET  /cctv-refresh    return fresh URL on HLS NETWORK_ERROR
  GET  /runtime-config  return current capture settings
  GET/POST/DELETE /roi, /calibration, /health
  WS /ws               subscribe to JSON broadcast
  WS /ws/detect        receive JPEG → inference → return annotated JPEG""",
        M, t+0.1, CW, 22.0, size=17)
    t += 22.3

    # ═══════════════════════════════════════════
    # RIGHT COLUMN
    # ═══════════════════════════════════════════
    t = CT

    # 5. Frontend Components
    t = sec_hdr(slide, "FRONTEND COMPONENTS", CR, t, CW)
    txtbox(slide, """[ App.jsx ]
  Key state: selectedCctv, cctvList, frameData, switching, calMode, pendingGps
  Core flow:
    Camera click → handleCctvClick (300ms debounce) → FlyTo animation
                 → POST /switch-camera → camera_ready WS signal → switching=false
    Calibration done → handleCalibSaved(heading) → update selectedCctv.heading
                     → update cctvList → MapView FOV direction auto-updated

[ MapView.jsx — deck.gl layer composition ]
  cctvHitLayer   (ScatterplotLayer, transparent)  click hit area
  cctvIconLayer  (IconLayer, SVG)                 CCTV icon (cyan when selected)
  cctvLabelLayer (TextLayer)                      camera name
  fovLayer       (PolygonLayer)                   selected camera FOV trapezoid
                                                  (nearM=15, farM=90, FOV=70 deg)
  trailLayer     (PathLayer)                      vehicle movement trail
  scatterLayer   (ScatterplotLayer)               vehicle position marker (shown at zoom>=15)
  textLayer      (TextLayer)                      vehicle track_id
  FOV trapezoid: polygon 15m-90m in heading direction, +-35 deg angle
                 auto-rotates to actual camera direction after calibration

[ CctvPlayer.jsx ]
  Tabs:
    Live  HLS stream playback (hls.js)
    YOLO  canvas capture → WS /ws/detect → display annotated image
    ROI   RoiEditor overlay
    Cal   CalibrationMode overlay
  HLS handling:
    - hls.destroy() + video.src="" + video.load() → fully clear previous frame
    - 15s timeout: force-release loading state if no MANIFEST_PARSED
    - NETWORK_ERROR → /cctv-refresh → restart with fresh URL
  Capture pipeline:
    - captureAndSend runs every captureIntervalMs (default 33ms)
    - maxInFlight (default 2): limit concurrent frames in-flight (preserves order)
    - JPEG quality 0.92, max width 640px

[ RoiEditor.jsx ]
  - Click: add vertex / Double-click: close polygon
  - Green semi-transparent fill
  - On save: POST /roi → stored in roi_manager
  - Auto-applied to detector on next camera switch

[ CalibrationMode.jsx ]
  Steps:
    1. Click road point on video (collect pixel coords)
    2. Instruction banner: "Click the same point on the map"
    3. Click on map (collect GPS coords)
    4. Repeat 1-3 for 4 pairs total
    5. POST /calibration → recompute homography + save calibration_data.json
       → compute bearing from GPS point 0 to 3 → onSaved(heading)
       → App updates heading → FOV direction auto-updated

[ useWebSocket.js ]
  Message routing:
    "camera_ready"  → increment cameraReady counter (clear switching state)
    "camera_error"  → show error message + clear switching
    others          → FrameAnalytics JSON → update frameData
  Auto-reconnect after 3s on disconnect""",
        CR, t+0.1, CW, 13.0, size=18)
    t += 13.3

    # 6. Key Workflows
    t = sec_hdr(slide, "KEY WORKFLOWS", CR, t, CW)
    txtbox(slide, """[ Workflow 1: Select CCTV and Start Detection ]
  User: click CCTV icon on map
    → App.handleCctvClick (300ms debounce)
    → setSelectedCctv(cctv), FlyTo animation
    → POST /switch-camera {cctvurl, lat, lon, name}
    → Backend: detector.reset_tracker(), analytics.reset()
               live_loop: stream.switch_to(new_url)
               WS broadcast: {type: "camera_ready"}
    → Frontend: switching=false, CctvPlayer starts HLS playback
    → User: click [YOLO tab] → vehicle detection displayed

[ Workflow 2: Set ROI ]
  User: CctvPlayer → [ROI tab]
    → RoiEditor canvas overlay shown
    → Draw polygon directly on video frame → [Save]
    → POST /roi {cctvurl, polygon}
    → roi_manager.save_roi(), detector.set_roi() applied immediately
    → Subsequent detections: objects outside ROI not passed to tracker

[ Workflow 3: Camera Direction Calibration ]
  User: CctvPlayer → [Calibration tab]
    → Click 4 road points on video (each paired with same point on map)
    → [Save] → POST /calibration
    → findHomography() → save calibration_data.json + swap transform matrix immediately
    → onSaved(heading) → update camera.heading → FOV trapezoid rotates
    → Auto-switch to ROI tab (recommended to set ROI right after calibration)

[ Workflow 4: HLS Token Expiry Recovery ]
  hls.js detects NETWORK_ERROR
    → GET /cctv-refresh?name=...&lat=...&lon=...
    → Re-call ITS API → return fresh HLS URL
    → hls.loadSource(newUrl) → stream restarts (no user action needed)""",
        CR, t+0.1, CW, 9.5, size=19)
    t += 9.8

    # 7. AI Detection & Tracking Pipeline
    t = sec_hdr(slide, "AI DETECTION & TRACKING PIPELINE", CR, t, CW)
    txtbox(slide, """[ YOLO Models ]
  yolov8x.engine  TensorRT FP16, ~9.6ms/frame (RTX 4070 baseline)
  yolov8s.engine  TensorRT FP16, ~3ms/frame
  yolov8x.pt      PyTorch GPU, ~33ms/frame
  * TensorRT .engine must be converted directly on the target GPU

[ BoxMOT Tracker Tiers ]
  auto    Auto-detect VRAM size → select appropriate tier
  cpu     ByteTrack  — no ReID, no VRAM needed, fastest
  low     OcSort     — no ReID, strong occlusion handling
  medium  BotSort    — appearance ReID, 6-8GB VRAM recommended
  high    DeepOcSort — appearance ReID, 8GB+ VRAM recommended
  ReID: re-assigns same track_id by appearance even after brief disappearance

[ Speed Calculation Logic ]
  1. Pixel coords → GPS coords (homography transform)
  2. Haversine(prev GPS, curr GPS) → dist_m
  3. dt = curr_timestamp - prev_timestamp (seconds)
  4. raw_kph = dist_m / dt x 3.6
  5. dist_m < 0.20 → speed=0 immediately (jitter removal)
  6. raw_kph > 180 → skip this frame, keep previous EMA (outlier removal)
  7. EMA: speed = 0.30 x raw + 0.70 x prev_ema
     (alpha=0.30: underestimated speeds corrected quickly)
  GC grace period: delete track only after 5+ undetected frames → speed continuity

[ Parked Vehicle Detection ]
  Method 1 (dwell): stationary 300 consecutive frames (~10s) → is_parked=True, register position
  Method 2 (position): new vehicle within 30px of registered position → is_parked immediately regardless of track_id
  Parked vehicles: excluded from stats/alerts, shown as grey marker on map""",
        CR, t+0.1, CW, 9.5, size=18)
    t += 9.8

    # 8. Performance
    t = sec_hdr(slide, "PERFORMANCE (RTX 4070 Laptop)", CR, t, CW)
    txtbox(slide, """[ Pipeline Latency Breakdown ]
  Canvas capture + JPEG encode    ~8ms
  WebSocket round-trip (localhost) ~3ms
  JPEG decode (server)             ~3ms
  YOLO + BoxMOT inference          ~9.6ms (yolov8x) / ~3ms (yolov8s)
  Annotation + JPEG encode         ~3ms
  ─────────────────────────────────────
  Total pipeline                   ~27ms (~37fps) yolov8x
                                   ~20ms (~50fps) yolov8s

[ TensorRT FP16 Conversion Effect ]
  CPU (yolov8x)              ~500ms/frame
  GPU + PyTorch (yolov8x)    ~33ms/frame
  TensorRT FP16 (yolov8x)    ~9.6ms/frame
  No accuracy change (FP16 error < 1%)

[ YOLO_DETECT_INTERVAL=3 Effect ]
  Every frame:      ~9.6ms x 30fps = 288ms/s GPU occupancy
  Every 3rd frame:  ~9.6ms x 10fps = 96ms/s GPU occupancy
  Other 2 frames: Kalman prediction only (~0.5ms) → more GPU headroom

[ model.predict(half=True) ]
  FP32 → FP16: ~30-50% inference speedup (.pt models only)
  TensorRT .engine already FP16, no change

[ inFlight Control (maxInFlight=2) ]
  Limit 2 concurrent frames → maximize throughput + preserve ordering""",
        CR, t+0.1, CW, 7.5, size=19)
    t += 7.8

    # 9. Design Decisions & Notes
    t = sec_hdr(slide, "DESIGN DECISIONS & NOTES", CR, t, CW)
    txtbox(slide, """[ No Shared Tracker Instance ]
  If live_loop and /ws/detect simultaneously update the same BoxMOT tracker,
  different frame sequences get mixed, completely breaking tracking.
  → Skip live_loop's track() call while /ws/detect is active.
    Call reset_tracker() when /ws/detect disconnects.

[ Auto-ROI Disabled ]
  Canny edge detection-based automatic ROI estimation mistook building/bridge edges
  for roads, filtering out most vehicles. Manual ROI only.

[ ITS API Single-Result Handling ]
  ITS API returns a dict (not a list) when there is only 1 result.
  → isinstance(raw, dict) branch applied to all API responses.

[ HLS Video Element Cleanup ]
  hls.destroy() only clears internal HLS state, not the video DOM buffer.
  → Must explicitly call video.src="" + video.load() on camera switch
    to prevent ghost frames from previous camera bleeding through.

[ 4-Point Homography Calibration ]
  Pixel→GPS homography implicitly encodes camera height, angle, and lens distortion.
  Calibration with 4 real road points is essential for accurate speed measurement.
  Default grid values are approximate and produce large errors.

[ Distance-Dependent Speed Error ]
  4-point homography is accurate only within the calibration quadrilateral.
  Vehicles at the top of the frame (farther away) are in the extrapolation zone
  → speed tends to be underestimated.
  Mitigation:
    1. update_from_calibration() updates _H_gps + _H_meter together (consistency)
    2. SPEED_SMOOTHING_ALPHA=0.30 — underestimated values corrected quickly
    3. Place 2 near + 2 far calibration points → better depth range coverage

[ useEffect Cleanup Scope ]
  Variables declared with const inside an if block within useEffect
  are inaccessible in the cleanup function outside that block (ReferenceError).
  → Declare variables needed in cleanup with let at the top of useEffect.

[ Tkinter Forward Reference ]
  _on_cuda_toggle in model_setup.py references fps_labels,
  which is populated on subsequent lines.
  Python closures use late binding (looked up at call time), so safe
  since the callback only fires after UI rendering is complete.""",
        CR, t+0.1, CW, 12.0, size=17)
    t += 12.3

    # 10. Demo screenshots
    t = sec_hdr(slide, "DEMO", CR, t, CW)
    t += 0.2
    img_h = 4.5
    img_gap = 0.4
    img_placeholder(slide, "Map view: vehicle dots, FOV trapezoid, speed alerts",
                    CR, t, CW, img_h)
    t += img_h + img_gap
    img_placeholder(slide, "CCTV player: YOLOv8 bounding boxes + track IDs",
                    CR, t, CW, img_h)
    t += img_h + img_gap
    img_placeholder(slide, "Calibration tab / ROI polygon editor",
                    CR, t, CW, img_h)

    prs.save("Traffic_Digital_Twin_Poster.pptx")
    print("Saved: Traffic_Digital_Twin_Poster.pptx")

if __name__ == "__main__":
    build()
